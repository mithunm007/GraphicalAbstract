import os
import json
from pptx import Presentation
from PyPDF2 import PdfReader
import requests
from jinja2 import Template
import re
import base64 # For encoding mock SVGs if needed, though direct embedding is easier
import io # Added for Streamlit file handling
import streamlit as st # Ensure streamlit is imported for the UI part

# --- Configuration for Local LLM ---
LOCAL_LLM_API_ENDPOINT = "http://localhost:11434/api/chat"
LOCAL_LLM_MODEL_NAME = "command-r7b" # Updated to deepseek-v2:latest as requested
API_TIMEOUT_SECONDS = 900
# --- End Configuration ---

# --- GLOBAL ICON LIBRARY ---
# These are professional-looking SVG icons that the LLM can explicitly reference.
ICON_LIBRARY = {
    "robot_ai": '''<svg class="icon text-blue-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M12 2C6.477 2 2 6.477 2 12s4.477 10 10 10 10-4.477 10-10S17.523 2 12 2zm0 2a8 8 0 110 16 8 8 0 010-16zm-3 5a1 1 0 100 2 1 1 0 000-2zm6 0a1 1 0 100 2 1 1 0 000-2zm-3 4a4 4 0 00-4 4h8a4 4 0 00-4-4z"/></svg>''',
    "cloud_platform": '''<svg class="icon text-green-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M18.73 11.27a6.5 6.5 0 00-11.46 0C5.1 11.96 4 13.56 4 15.33A5.5 5.5 0 009.5 20h7A5.5 5.5 0 0022 14.5c0-2.3-1.6-4.32-3.27-5.23zM9.5 18A3.5 3.5 0 016 14.5c0-1.2.66-2.27 1.7-2.88L8 11.27l.73-.27C9.28 10.5 9.77 10 10.33 10a4.5 4.5 0 018.17 2.37L19 12.5l.38.12A3.5 3.5 0 0122 14.5c0 1.93-1.57 3.5-3.5 3.5h-7z"/></svg>''',
    "data_analytics": '''<svg class="icon text-teal-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M16 11V3h-2v8h-4V7H8v4H4v2h4v4h2v-4h4v4h2v-6zM18 19h-2v-2h2v2zm0-4h-2v-2h2v2z"/></svg>''',
    "gear_solution": '''<svg class="icon text-purple-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M12 2C6.477 2 2 6.477 2 12s4.477 10 10 10 10-4.477 10-10S17.523 2 12 2zm0 2a8 8 0 110 16 8 8 0 010-16zM11 7h2v6h-2V7zm0 8h2v2h-2v-2z"/></svg>''',
    "chart_growth": '''<svg class="icon text-orange-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M16 6l-4 4-4-4-6 6v2h16v-2l-4-4zM2 18h20v2H2z"/></svg>''',
    "product_box": '''<svg class="icon text-red-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M20 6h-4V4c0-1.1-.9-2-2-2h-4c-1.1 0-2 .9-2 2v2H4c-1.1 0-2 .9-2 2v12c0 1.1.9 2 2 2h16c1.1 0 2-.9 2-2V8c0-1.1-.9-2-2-2zm-6 0h-4V4h4v2z"/></svg>''',
    "document_report": '''<svg class="icon text-gray-700" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M14 2H6c-1.1 0-1.99.9-1.99 2L4 20c0 1.1.89 2 1.99 2H18c1.1 0 2-.9 2-2V8l-6-6zM6 20V4h7v5h5v11H6z"/></svg>''',
    "globe_network": '''<svg class="icon text-indigo-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M12 2C6.48 2 2 6.48 2 12s4.48 10 10 10 10-4.48 10-10S17.52 2 12 2zm-1 17.93c-3.95-.49-7-3.8-7-7.93 0-2.48 1.19-4.73 3.04-6.14L11 8.52v9.41zm2 0V8.52l4.96-6.14C20.81 7.27 22 9.52 22 12c0 4.13-3.05 7.44-7 7.93zm-8-9.41l-1.96-2.43C5.19 8.27 4 9.93 4 12c0 1.48.64 2.83 1.66 3.77L9 14.52zM19 14.52l1.96 2.43C18.81 19.73 17.36 21 15.7 21c-1.48 0-2.83-.64-3.77-1.66L13 9.48z"/></svg>''',
    "cog_settings": '''<svg class="icon text-gray-600" fill="currentColor" viewBox="0 0 24 24" xmlns="http://www.w3.org/2000/svg"><path d="M12 2c-1.1 0-2 .9-2 2v3.1c-1.5-.4-3.1-.3-4.5.3l-2.4-1.4c-.9-.5-2.1-.2-2.5.8l-1.2 2.1c-.4.8-.2 1.8.6 2.3l2.4 1.4c-.1.5-.1 1-.1 1.5s0 1 .1 1.5l-2.4 1.4c-.8.5-1.1 1.5-.6 2.3l1.2 2.1c.4.9 1.6 1.2 2.5.8l2.4-1.4c1.4.6 3 .7 4.5.3V20c0 1.1.9 2 2 2h4c1.1 0 2-.9 2-2v-3.1c1.5.4 3.1.3 4.5-.3l2.4 1.4c.9.5 2.1.2 2.5-.8l1.2-2.1c.4-.8.2-1.8-.6-2.3l-2.4-1.4c.1-.5.1-1 .1-1.5s0-1-.1-1.5l2.4-1.4c.8-.5 1.1-1.5.6-2.3l-1.2-2.1c-.4-.9-1.6-1.2-2.5-.8l-2.4 1.4c-1.4-.6-3-.7-4.5-.3V4c0-1.1-.9-2-2-2h-4zM12 15c-1.66 0-3-1.34-3-3s1.34-3 3-3 3 1.34 3 3-1.34 3-3 3z"/></svg>'''
}

# --- MOCKED EXTERNAL TOOL FUNCTIONS ---
def mock_call_svg_dreamer_or_omnisvg(prompt: str) -> str:
    """
    Mocks a call to an SVG generation model (SVGDreamer/OmniSVG).
    Prioritizes specific icon names suggested by LLM, then keyword matching, then a generic default.
    """
    print(f"MOCK CALL: Generating SVG for prompt: '{prompt}'")
    
    prompt_lower = prompt.lower()

    # 1. Check if the prompt is an exact match for a predefined icon name
    if prompt in ICON_LIBRARY:
        return ICON_LIBRARY[prompt]

    # 2. Check for keywords and return a relevant icon from the library
    if "ai" in prompt_lower or "agent" in prompt_lower or "intelligence" in prompt_lower or "robot" in prompt_lower:
        return ICON_LIBRARY["robot_ai"]
    elif "platform" in prompt_lower or "cloud" in prompt_lower or "network" in prompt_lower or "system" in prompt_lower:
        return ICON_LIBRARY["cloud_platform"]
    elif "data" in prompt_lower or "analytics" in prompt_lower or "insights" in prompt_lower or "chart" in prompt_lower:
        return ICON_LIBRARY["data_analytics"]
    elif "solution" in prompt_lower or "service" in prompt_lower or "offerings" in prompt_lower or "development" in prompt_lower or "api" in prompt_lower:
        return ICON_LIBRARY["gear_solution"]
    elif "market" in prompt_lower or "analysis" in prompt_lower or "sales" in prompt_lower or "trend" in prompt_lower or "growth" in prompt_lower:
        return ICON_LIBRARY["chart_growth"]
    elif "product" in prompt_lower or "overview" in prompt_lower or "features" in prompt_lower or "specifications" in prompt_lower:
        return ICON_LIBRARY["product_box"]
    elif "report" in prompt_lower or "document" in prompt_lower or "summary" in prompt_lower or "section" in prompt_lower:
        return ICON_LIBRARY["document_report"]
    elif "globe" in prompt_lower or "global" in prompt_lower:
        return ICON_LIBRARY["globe_network"]
    elif "settings" in prompt_lower or "configure" in prompt_lower:
        return ICON_LIBRARY["cog_settings"]
    
    # 3. Fallback to a generic professional default icon if no specific keyword matches
    return ICON_LIBRARY["document_report"] # Using a generic document icon as a slightly better default

def mock_call_chart_gpt(data_summary: dict) -> dict:
    """
    Mocks a call to a chart generation model (ChartGPT).
    Returns structured data suitable for Chart.js with varied chart types.
    """
    print(f"MOCK CALL: Generating chart data for summary: {data_summary}")
    labels = [m['label'] for m in data_summary.get('metrics', [])]
    values = [float(re.sub(r'[^\d.]', '', str(m['value']))) if re.sub(r'[^\d.]', '', str(m['value'])) else 0 for m in data_summary.get('metrics', [])]

    chart_type = data_summary.get('chart_type', 'bar').lower()

    base_colors = ['#63b3ed', '#a78bfa', '#81e6d9', '#f6ad55', '#fc8181', '#718096']
    background_colors = [base_colors[i % len(base_colors)] for i in range(len(labels))]
    border_colors = [c.replace('3', '2').replace('6', '5') for c in background_colors] # Slightly darker for border

    datasets = [{
        "label": data_summary.get('chart_title', "Metrics"),
        "data": values,
        "backgroundColor": background_colors,
        "borderColor": border_colors,
        "borderWidth": 1
    }]

    if chart_type == 'line':
        datasets[0]["fill"] = False
        datasets[0]["tension"] = 0.3 # Smooth lines
        datasets[0]["borderColor"] = base_colors[0] # Single color for line
        datasets[0]["backgroundColor"] = base_colors[0]
        datasets[0]["pointRadius"] = 5
        datasets[0]["pointHoverRadius"] = 7
    elif chart_type == 'doughnut' or chart_type == 'pie':
        # For pie/doughnut, backgroundColors apply to segments
        pass # Colors are already set for segments

    return {
        "type": chart_type,
        "data": {
            "labels": labels,
            "datasets": datasets
        },
        "options": {
            "responsive": True,
            "maintainAspectRatio": False,
            "plugins": {
                "title": {
                    "display": True,
                    "text": data_summary.get('chart_title', "Metrics Overview"),
                    "font": {
                        "size": 16,
                        "weight": 'bold'
                    },
                    "color": '#2d3748'
                },
                "legend": {
                    "display": True,
                    "position": 'bottom',
                    "labels": {
                        "color": '#4a5568'
                    }
                }
            },
            "scales": {
                "y": {
                    "beginAtZero": True,
                    "display": chart_type not in ['pie', 'doughnut'], # Hide Y-axis for pie/doughnut
                    "ticks": {
                        "color": '#4a5568'
                    },
                    "grid": {
                        "color": 'rgba(200, 200, 200, 0.2)'
                    }
                },
                "x": {
                    "display": chart_type not in ['pie', 'doughnut'], # Hide X-axis for pie/doughnut
                    "ticks": {
                        "color": '#4a5568'
                    },
                    "grid": {
                        "color": 'rgba(200, 200, 200, 0.2)'
                    }
                }
            },
            "animation": { # Add some animation for professional look
                "duration": 1000,
                "easing": 'easeOutQuart'
            }
        }
    }

def mock_call_deplot(image_path: str) -> dict:
    """
    Mocks a call to DePlot. In a real scenario, this would read an image,
    send it to DePlot, and return the extracted tabular data.
    """
    print(f"MOCK CALL: Extracting data from chart image: '{image_path}' (This is a placeholder)")
    # Placeholder for extracted data
    return {
        "extracted_data": [
            {"year": 2020, "value": 100},
            {"year": 2021, "value": 120},
            {"year": 2022, "value": 150}
        ],
        "notes": "Data extracted from mock image. Actual DePlot output would be more precise."
    }

# --- End MOCKED EXTERNAL TOOL FUNCTIONS ---


# Function to call the local LLM API (Ollama)
def call_llm_api(prompt: str) -> str:
    print(f"\n--- Attempting to call Local LLM API ({LOCAL_LLM_MODEL_NAME}) ---")
    print("Prompt (excerpt):\n" + prompt[:500] + "...")
    print("---------------------------------------------")

    headers = {"Content-Type": "application/json"}
    payload = {
        "model": LOCAL_LLM_MODEL_NAME,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "stream": False,
        "options": {
            "temperature": 0.2, # Very low temperature for structured output
            "num_gpu": -1 # Instruct Ollama to use all available GPU layers
        }
    }

    try:
        response = requests.post(
            LOCAL_LLM_API_ENDPOINT, headers=headers, json=payload, timeout=API_TIMEOUT_SECONDS
        )
        response.raise_for_status() # Raise an HTTPError for bad responses (4xx or 5xx)
        
        generated_text = ""
        response_json = response.json()

        # Extract content from various possible Ollama/LLM response structures
        if "message" in response_json and "content" in response_json["message"]:
            generated_text = response_json["message"]["content"]
        elif "response" in response_json: # For some Ollama raw API outputs
            generated_text = response_json.get("response", "")
        elif "choices" in response_json and response_json["choices"] and \
             "message" in response_json["choices"][0] and "content" in response_json["choices"][0]["message"]:
            generated_text = response_json["choices"][0]["message"]["content"]

        if not generated_text:
            print(f"Warning: No text found in LLM response. Raw response: {response.text}")
            return json.dumps({"error": "LLM response was empty or malformed."})

        # --- IMPORTANT: Improved JSON extraction logic ---
        # First, try to parse the entire generated_text as JSON
        try:
            parsed_json = json.loads(generated_text)
            # If successful, re-serialize it to a string for consistent return,
            # ensuring proper formatting and removal of extra whitespace/newlines.
            extracted_json_str = json.dumps(parsed_json, indent=None, separators=(',', ':'))
            print("Info: LLM returned raw JSON (successfully parsed).")
        except json.JSONDecodeError:
            # If direct parsing fails, try to find JSON within a markdown block
            json_match = re.search(r'```json\n(.*?)\n```', generated_text, re.DOTALL)
            if json_match:
                extracted_json_str = json_match.group(1).strip()
                # Remove common trailing commas that might cause JSONDecodeError (often from LLMs)
                extracted_json_str = re.sub(r',\s*}', '}', extracted_json_str)
                extracted_json_str = re.sub(r',\s*]', ']', extracted_json_str)
                print("Info: LLM returned JSON in markdown block (successfully parsed).")
            else:
                print(f"Warning: No JSON block found and raw output is not valid JSON. Raw output:\n{generated_text}")
                return json.dumps({"error": "LLM did not return valid JSON or JSON in expected markdown block."})

        if extracted_json_str:
            print("\n--- Extracted JSON from LLM response (preview) ---")
            # print(extracted_json_str) # Uncomment for full JSON in console for debugging
            print("------------------------------------------")
            return extracted_json_str
        else:
            # This case should ideally not be reached with the improved logic above
            return json.dumps({"error": "Failed to extract any JSON from LLM response."})

    except requests.exceptions.RequestException as e:
        error_message = f"Error communicating with local LLM API at {LOCAL_LLM_API_ENDPOINT}: {e}"
        print(error_message)
        return json.dumps({"error": error_message})
    except (KeyError, json.JSONDecodeError) as e:
        # This block now primarily catches errors after initial extraction attempts
        error_message = f"Error parsing local LLM API response: {e}. Raw response: {response.text if 'response' in locals() else 'No response object'}"
        print(error_message)
        return json.dumps({"error": error_message})
    except TypeError as e:
        print(f"Type error during LLM API call: {e}. Raw output: {generated_text if 'generated_text' in locals() else 'No generated text'}")
        return json.dumps({"error": f"Type error in LLM response: {e}"})


# Functions to extract text from different document types
# These functions will now receive file-like objects (BytesIO) from Streamlit's uploaded_file
def extract_text_from_txt(uploaded_file_stream) -> str:
    """Extracts text from a .txt file stream."""
    return uploaded_file_stream.getvalue().decode('utf-8')

def extract_text_from_pdf(uploaded_file_stream) -> str:
    """Extracts text from a .pdf file stream."""
    try:
        text = ""
        reader = PdfReader(uploaded_file_stream)
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""

def extract_text_from_pptx(uploaded_file_stream) -> str:
    """Extracts text from a .pptx file stream."""
    try:
        prs = Presentation(uploaded_file_stream)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""

# Add docx and xml/html extraction functions for Streamlit's uploaded_file
try:
    from docx import Document # For .docx files
except ImportError:
    Document = None
    # st.warning is now in Streamlit UI section if needed

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
    # st.warning is now in Streamlit UI section if needed

def extract_text_from_docx(uploaded_file_stream) -> str:
    """Extracts text from a .docx file uploaded via Streamlit."""
    if Document is None:
        # This warning is now handled in the Streamlit UI setup
        return ""
    try:
        document = Document(uploaded_file_stream)
        full_text = [paragraph.text for paragraph in document.paragraphs]
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_xml(uploaded_file_stream) -> str:
    """Extracts text content from an XML/HTML file uploaded via Streamlit."""
    if BeautifulSoup is None:
        # This warning is now handled in the Streamlit UI setup
        xml_content = uploaded_file_stream.getvalue().decode('utf-8')
        clean_text = re.sub(r'<[^>]+>', '', xml_content)
        return clean_text
        
    try:
        xml_content = uploaded_file_stream.getvalue().decode('utf-8')
        # Attempt to parse as XML first, then as HTML if XML parsing fails
        try:
            soup = BeautifulSoup(xml_content, 'xml') # Try XML parser
            text = soup.get_text(separator='\n', strip=True)
            if not text.strip(): # If XML parser found no text, it might be HTML content
                raise ValueError("XML parsing yielded no meaningful text, trying HTML parser.")
            return text
        except Exception as xml_e: # Catch all exceptions from XML parsing
            print(f"XML parsing failed: {xml_e}. Attempting HTML parsing.")
            soup = BeautifulSoup(xml_content, 'html.parser') # Fallback to HTML parser
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        print(f"Error extracting text from XML/HTML file: {e}. Ensure the file is well-formed.")
        return ""


# Main function to dispatch document text extraction for Streamlit uploaded files
def extract_text_from_document(uploaded_file) -> str:
    file_type = uploaded_file.type
    # Streamlit uploaded files provide a BytesIO object directly for content
    # Ensure the stream is reset to the beginning before reading
    uploaded_file.seek(0)
    
    if "text/plain" in file_type: # .txt
        return extract_text_from_txt(uploaded_file)
    elif "application/pdf" in file_type: # .pdf
        return extract_text_from_pdf(uploaded_file)
    elif "application/vnd.openxmlformats-officedocument.presentationml.presentation" in file_type: # .pptx
        return extract_text_from_pptx(uploaded_file)
    elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in file_type: # .docx
        return extract_text_from_docx(uploaded_file)
    elif "text/xml" in file_type or "application/xml" in file_type or "text/html" in file_type: # .xml or .html
        return extract_text_from_xml(uploaded_file)
    else:
        # This warning is now handled in the Streamlit UI setup
        return ""

# Placeholder for image extraction from PDF/PPTX - this is a complex task
def extract_images_from_document_mock(uploaded_file) -> list[str]:
    """
    Mocks image extraction. In a real scenario, this would parse the document
    and return paths to extracted image files (e.g., charts).
    """
    print(f"MOCK: Attempting to extract images from {uploaded_file.name}. (Not implemented in full for this example).")
    # Simulate finding a chart image based on file name or type
    if "chart" in uploaded_file.name.lower() or "report" in uploaded_file.name.lower() or "pdf" in uploaded_file.type:
        return ["mock_chart_image_from_doc.png"] # Return a dummy path
    return []


# Main function to generate the graphical abstract (returns HTML string)
def generate_graphical_abstract_html(document_content: str) -> str:
    """
    Generates the graphical abstract HTML based on document content.
    Returns the HTML string.
    """
    # Define the NEW JSON schema template including secondary metrics section
    json_schema_template = """{
    "main_title": "string",
    "about_section": {
        "title": "string",
        "summary": "string"
    },
    "metrics_section": {
        "title": "string",
        "metrics": [
            {"value": "string", "label": "string"}
        ],
        "chart_type": "string",  // e.g., "bar", "pie", "line", "doughnut"
        "chart_title": "string",
        "chart_description": "string" // LLM to provide textual summary for ChartGPT
    },
    "secondary_metrics_section": { // NEW SECTION for a second graph
        "title": "string",
        "metrics": [
            {"value": "string", "label": "string"}
        ],
        "chart_type": "string",  // e.g., "bar", "pie", "line", "doughnut"
        "chart_title": "string",
        "chart_description": "string"
    },
    "offerings_section": {
        "title": "string",
        "items": [
            {"title": "string", "description": "string", "illustration_prompt": "string", "suggested_icon_name": "string"} // Added suggested_icon_name
        ]
    },
    "conclusion_section": {
        "title": "string",
        "understanding_section": {
            "title": "string",
            "description": "string",
            "key_concepts": ["string"]
        },
        "outlook_section": {
            "title": "string",
            "info_box_text": "string",
            "description": "string"
        }
    },
    "insights_and_implications_section": {
        "title": "string",
        "primary_insights": {
            "title": "string",
            "description": "string"
        },
        "practical_implications": {
            "title": "string",
            "description": "string"
        },
        "broader_significance": {
            "title": "string",
            "description": "string"
        }
    },
    "footer_text": "string"
}"""

    # Get available icon names for LLM to reference
    available_icon_names = ", ".join(f"'{name}'" for name in ICON_LIBRARY.keys())


    # Construct the prompt for the LLM with updated instructions for dynamic headings and content length
    llm_prompt = f"""
You are an expert at extracting COMPREHENSIVE, DETAILED, and RELEVANT information from documents and structuring it into JSON FORMAT for a GRAPHICAL ABSTRACT GENERATION. Your goal is to make the summaries and descriptions as informative and extensive as possible, reflecting the depth of the original document. Crucially, **ensure that every single field in the provided JSON schema is filled with content**. If the document does not explicitly provide information for a field, infer plausible and contextually relevant content based on the overall document theme to ensure no blank spaces or "string" placeholders appear in the final output. Aim for at least 2-3 sentences for each summary/description field.

Analyze the following document content and extract information to accurately fill the JSON structure below.
For each field, generate a title or description that is **dynamically relevant to the content extracted from the document**, rather than using generic placeholders. Aim for clear and informative summaries.

**Specifically for 'illustration_prompt' in 'offerings_section.items'**: Provide a short, descriptive phrase that could be used to generate a relevant vector illustration (e.g., "a secure network diagram", "a fast-moving gear").
**Specifically for 'suggested_icon_name' in 'offerings_section.items'**: From the list [{available_icon_names}], choose the most appropriate icon name for the item. This is crucial for selecting correct visuals. If no perfect match, choose a conceptually related one.

**Specifically for 'metrics_section' and 'secondary_metrics_section'**:
- Populate 'metrics' array with actual values and labels from the document.
- Based on the nature of the metrics, suggest a suitable 'chart_type' (e.g., "bar", "pie", "line", "doughnut"). Try to vary the chart types for different data sets to make them more engaging.
- Provide a 'chart_title' and a 'chart_description' summarizing the data for the chart.
- **CRITICAL**: Ensure you generate data for at least two distinct charts (one for 'metrics_section' and one for 'secondary_metrics_section') and a total of at least 4 numerical metrics across both sections. If not enough *real, explicit numerical data* for 4 metrics, invent plausible, thematic dummy values (e.g., "75%", "12.3M", "Q3 2024 Growth") with relevant labels to reach the minimum of 4 metrics. These invented metrics must directly relate to the overall document topic.

**For 'offerings_section.items'**: Extract relevant offerings from the document. **CRITICAL: Ensure exactly 4 distinct, plausible, relevant, and engaging items are generated for this section, ALWAYS.** If the document provides fewer than 4 distinct items, invent additional ones that strongly fit the document's overall theme and content. These invented items should have descriptive but generic titles like "Strategic Consulting" or "Innovative Solutions" and should NOT use phrases like 'Coming Soon'. Ensure each item (whether extracted or invented) has a 'title', 'description', 'illustration_prompt', and a 'suggested_icon_name' chosen from the provided `ICON_LIBRARY` that is colorful and relevant (e.g., 'robot_ai', 'cloud_platform', 'data_analytics').

**For all other sections (e.g., 'conclusion_section', 'insights_and_implications_section')**: If the document does not explicitly provide content for these sections, **generate a concise summary and relevant details based on the overall themes and information present in the document**. Do not leave any section title or description blank, use "string" placeholders, or omit the section entirely if it's part of the schema.

**CRITICAL INSTRUCTION FOR DOCUMENT CONTENT PROCESSING**: When analyzing the "Document Content to Analyze" below, **explicitly ignore and do not process any sections that appear to be "References," "Bibliography," "Citations," "Acknowledgements," or similar academic/research listing sections.** Focus solely on the core content of the paper for abstract generation.

**JSON Schema to adhere to (STRICTLY):**
```json
{json_schema_template}
```

**CRITICAL INSTRUCTION: Your response MUST contain ONLY the JSON output, enclosed within a single markdown code block.
Do NOT include any other text, explanations, or conversational filler outside this JSON block.**
Ensure all string values in the JSON are plain text, without Markdown or HTML.

**Document Content to Analyze:**
---
{document_content}
---

Generate ONLY the JSON output according to the schema and strict content rules.
"""

    # Call the LLM to generate the JSON data
    st.info("Calling local LLM to generate structured data...")
    json_output_text = call_llm_api(llm_prompt)

    try:
        abstract_data = json.loads(json_output_text)
        # --- NEW: Robust type checking for abstract_data ---
        if isinstance(abstract_data, list):
            if abstract_data: # If it's a non-empty list, try to get the first item if it's a dict
                if isinstance(abstract_data[0], dict):
                    abstract_data = abstract_data[0]
                    print("Info: LLM returned a JSON list, but the first element is a dictionary. Using it.")
                else:
                    st.error(f"Error: LLM returned a JSON list, and its first element is not a dictionary. Cannot process. Raw output: {json_output_text}")
                    return None
            else:
                st.error(f"Error: LLM returned an empty JSON list. Cannot process. Raw output: {json_output_text}")
                return None
        elif not isinstance(abstract_data, dict):
            st.error(f"Error: LLM returned unexpected JSON type ({type(abstract_data).__name__}). Expected a dictionary. Raw output: {json_output_text}")
            return None
        # --- END NEW: Robust type checking ---

        if "error" in abstract_data: # Check if the LLM call itself returned an error JSON
            st.error(f"Error from LLM (returned JSON with error): {abstract_data['error']}")
            return None
    except json.JSONDecodeError as e:
        st.error(f"Error parsing JSON from LLM: {e}")
        st.code(f"LLM Raw Output (might be incomplete/invalid JSON): {json_output_text}")
        return None
    except TypeError as e:
        st.error(f"Type error during LLM API call: {e}. Raw output: {json_output_text}")
        return None

    # DEBUG: Print raw abstract_data directly after JSON load
    print(f"DEBUG: Raw abstract_data after JSON load:\n{json.dumps(abstract_data, indent=2)}")

    # --- Post-processing: Add SVG icons and clean up any LLM placeholders ---
    def clean_and_get(data, key, default_value="", is_list=False):
        # This function ensures that if a key is missing or contains "string", it returns a default value.
        # It's crucial for pre-processing data before passing to Jinja.
        if data is None:
            return default_value if not is_list else []
        if is_list:
            val = data.get(key, [])
            cleaned_list = []
            for item in val:
                if isinstance(item, dict):
                    cleaned_item_dict = {}
                    for k, v in item.items():
                        if isinstance(v, str) and v.strip() != "string":
                            cleaned_item_dict[k] = v.strip()
                        elif not isinstance(v, str):
                            cleaned_item_dict[k] = v
                    if cleaned_item_dict:
                        cleaned_list.append(cleaned_item_dict)
                elif isinstance(item, str) and item.strip() != "" and item.strip() != "string":
                     cleaned_list.append(item.strip())
            return [s for s in cleaned_list if s]
        else:
            val = data.get(key, default_value)
            return val.replace("string", "").strip() if isinstance(val, str) else val

    # Assign parsed and cleaned data to abstract_data structure
    # The LLM is now instructed to always provide a main title, so default can be more generic or empty if LLM ensures content.
    abstract_data['main_title'] = clean_and_get(abstract_data, 'main_title', 'Graphical Abstract Summary')

    # Ensure all top-level sections are dictionaries to prevent UndefinedError in Jinja
    abstract_data['about_section'] = abstract_data.get('about_section', {})
    abstract_data['metrics_section'] = abstract_data.get('metrics_section', {})
    abstract_data['secondary_metrics_section'] = abstract_data.get('secondary_metrics_section', {})
    abstract_data['offerings_section'] = abstract_data.get('offerings_section', {})
    abstract_data['conclusion_section'] = abstract_data.get('conclusion_section', {})
    abstract_data['insights_and_implications_section'] = abstract_data.get('insights_and_implications_section', {})

    # Populate all sub-sections and their fields with content or sensible defaults/inferences
    # About Section
    about_sec = abstract_data['about_section']
    about_sec['title'] = clean_and_get(about_sec, 'title', 'About the Document')
    about_sec['summary'] = clean_and_get(about_sec, 'summary', 'This section provides a brief overview of the document\'s core content and purpose.')

    # Primary Metrics Section
    met_sec = abstract_data['metrics_section'] 
    met_sec['title'] = clean_and_get(met_sec, 'title', 'Key Metrics Summary')
    met_sec['chart_type'] = clean_and_get(met_sec, 'chart_type', 'bar').lower()
    if met_sec['chart_type'] not in ['bar', 'line', 'pie', 'doughnut']:
         met_sec['chart_type'] = 'bar'
    met_sec['chart_title'] = clean_and_get(met_sec, 'chart_title', 'Key Metrics Overview')
    met_sec['chart_description'] = clean_and_get(met_sec, 'chart_description', 'Visual representation of key organizational metrics or research findings.')
    
    # Secondary Metrics Section
    sec_met_sec = abstract_data['secondary_metrics_section'] 
    sec_met_sec['title'] = clean_and_get(sec_met_sec, 'title', 'Additional Insights')
    sec_met_sec['chart_type'] = clean_and_get(sec_met_sec, 'chart_type', 'line').lower()
    if sec_met_sec['chart_type'] not in ['bar', 'line', 'pie', 'doughnut']:
         sec_met_sec['chart_type'] = 'line'
    sec_met_sec['chart_title'] = clean_and_get(sec_met_sec, 'chart_title', 'Trends and Growth Analysis')
    sec_met_sec['chart_description'] = clean_and_get(sec_met_sec, 'chart_description', 'Detailed analysis of trends and growth patterns within the document\'s scope.')

    # Ensure metrics are numeric and cleaned
    def process_metrics(metrics_list, default_label_prefix):
        cleaned_metrics = []
        for i, m in enumerate(metrics_list):
            if not isinstance(m, dict) or 'value' not in m or 'label' not in m:
                print(f"Warning: Invalid metric item encountered: {m}. Skipping.")
                continue

            raw_value = m['value']
            label = clean_and_get(m, 'label')
            numeric_value = None

            if isinstance(raw_value, (int, float)):
                numeric_value = float(raw_value)
            elif isinstance(raw_value, str):
                num_match = re.search(r'[-+]?\d*\.?\d+', raw_value)
                if num_match:
                    try:
                        numeric_value = float(num_match.group(0))
                    except ValueError:
                        pass # numeric_value remains None
            
            # Ensure value is numeric and label is not empty/placeholder
            if numeric_value is not None and label and label.strip() != "" and label.strip() != "string":
                cleaned_metrics.append({'value': numeric_value, 'label': label.strip()})
            else:
                # Invent a plausible dummy metric if LLM didn't provide enough or valid ones
                cleaned_metrics.append({
                    'value': float(50 + (len(cleaned_metrics) * 5)), # Simple progressive dummy value
                    'label': f"{default_label_prefix} {i + 1}"
                })
        return cleaned_metrics

    met_sec['metrics'] = process_metrics(clean_and_get(met_sec, 'metrics', is_list=True), "Primary Metric")
    sec_met_sec['metrics'] = process_metrics(clean_and_get(sec_met_sec, 'metrics', is_list=True), "Secondary Metric")

    # Ensure there are at least 4 total metrics, supplementing with invented ones if needed.
    total_metrics_count = len(met_sec['metrics']) + len(sec_met_sec['metrics'])
    if total_metrics_count < 4:
        num_invented_needed = 4 - total_metrics_count
        for i in range(num_invented_needed):
            # Add to primary if it has less than 2, otherwise add to secondary
            if len(met_sec['metrics']) < 2:
                met_sec['metrics'].append({
                    'value': float(70 + i * 3), # Invented value
                    'label': f"Inferred Metric {len(met_sec['metrics']) + 1}"
                })
            else:
                 sec_met_sec['metrics'].append({
                    'value': float(40 + i * 3), # Invented value
                    'label': f"Inferred Secondary Metric {len(sec_met_sec['metrics']) + 1}"
                })


    # Chart.js data generation
    met_sec['chart_js_data'] = mock_call_chart_gpt(met_sec)
    sec_met_sec['chart_js_data'] = mock_call_chart_gpt(sec_met_sec)

    print("Final cleaned primary metrics:", met_sec['metrics'])
    print("Final cleaned secondary metrics:", sec_met_sec['metrics'])

    # Offerings Section - Ensure exactly 4 items, generating if necessary
    off_sec = abstract_data['offerings_section']
    off_sec['title'] = clean_and_get(off_sec, 'title', 'Our Core Offerings')
    
    processed_offerings = []
    # Filter out empty/placeholder items from LLM response first
    for item in clean_and_get(off_sec, 'items', is_list=True):
        if clean_and_get(item, 'title') and clean_and_get(item, 'description'):
            processed_offerings.append(item)
    
    # Ensure exactly 4 offerings are present
    current_offerings_count = len(processed_offerings)
    if current_offerings_count < 4:
        num_to_generate = 4 - current_offerings_count
        # Generic, plausible offerings to be generated if the LLM didn't provide enough
        generic_thematic_offerings = [
            {"title": "Strategic Consulting", "description": "Expert advisory on market trends, competitive landscapes, and growth strategies.", "suggested_icon_name": "chart_growth"},
            {"title": "Custom Solution Development", "description": "Building bespoke software and AI applications tailored to specific business challenges.", "suggested_icon_name": "gear_solution"},
            {"title": "Data Analytics & Insights", "description": "Transforming raw data into actionable insights for informed decision-making.", "suggested_icon_name": "data_analytics"},
            {"title": "Managed Services", "description": "Comprehensive support and maintenance for your technology infrastructure and platforms.", "suggested_icon_name": "cloud_platform"}
        ]
        
        for i in range(num_to_generate):
            # Cycle through generic_thematic_offerings, ensuring unique icons if possible
            idx = (current_offerings_count + i) % len(generic_thematic_offerings)
            new_item = generic_thematic_offerings[idx].copy() # Copy to avoid modifying original template
            processed_offerings.append(new_item)

    # Assign SVGs to all processed offerings
    for item in processed_offerings:
        suggested_icon_name = clean_and_get(item, 'suggested_icon_name', '').strip()
        if suggested_icon_name and suggested_icon_name in ICON_LIBRARY:
            item['illustration_svg'] = ICON_LIBRARY[suggested_icon_name]
        else:
            item['illustration_svg'] = mock_call_svg_dreamer_or_omnisvg(clean_and_get(item, 'illustration_prompt', item.get('title', 'concept')))
    
    off_sec['items'] = processed_offerings # Set the final list of 4 items

    # Conclusion Section
    conc_sec = abstract_data['conclusion_section']
    conc_sec['title'] = clean_and_get(conc_sec, 'title', 'Key Conclusions')
    
    conc_sec['understanding_section'] = conc_sec.get('understanding_section', {})
    ua = conc_sec['understanding_section']
    ua['title'] = clean_and_get(ua, 'title', 'Understanding the Core Message')
    ua['description'] = clean_and_get(ua, 'description', 'This section consolidates the most important findings and overarching message from the document.')
    ua['key_concepts'] = clean_and_get(ua, 'key_concepts', is_list=True)
    if not ua['key_concepts']: # Ensure at least some key concepts
        ua['key_concepts'] = ["Innovation", "Strategic Impact", "Future Growth"] # Generic but relevant
    
    conc_sec['outlook_section'] = conc_sec.get('outlook_section', {})
    ol = conc_sec['outlook_section']
    ol['title'] = clean_and_get(ol, 'title', 'Future Outlook')
    ol['info_box_text'] = clean_and_get(ol, 'info_box_text', 'Exploring potential future developments and their implications.')
    ol['description'] = clean_and_get(ol, 'description', 'The document\'s insights suggest a promising trajectory for future research and application.')


    # Insights and Implications Section
    imp_sec = abstract_data['insights_and_implications_section']
    imp_sec['title'] = clean_and_get(imp_sec, 'title', 'Insights and Implications')
    
    imp_sec['primary_insights'] = imp_sec.get('primary_insights', {})
    pi = imp_sec['primary_insights']
    pi['title'] = clean_and_get(pi, 'title', 'Primary Insights')
    pi['description'] = clean_and_get(pi, 'description', 'The most significant discoveries or understandings derived from the analysis.')

    imp_sec['practical_implications'] = imp_sec.get('practical_implications', {})
    pai = imp_sec['practical_implications']
    pai['title'] = clean_and_get(pai, 'title', 'Practical Implications')
    pai['description'] = clean_and_get(pai, 'description', 'How the findings can be applied in real-world scenarios or industry practices.')
    
    imp_sec['broader_significance'] = imp_sec.get('broader_significance', {})
    bs = imp_sec['broader_significance']
    bs['title'] = clean_and_get(bs, 'title', 'Broader Significance')
    bs['description'] = clean_and_get(bs, 'description', 'The wider impact and importance of the document\'s content on the field or society.')
    
    abstract_data['footer_text'] = clean_and_get(abstract_data, 'footer_text', '2025 Graphical Abstract. All rights reserved.') # Generic default for footer
    
    # --- DIAGNOSTIC PRINT: Check final abstract_data before Jinja2 rendering ---
    print("\n--- Final abstract_data before Jinja2 rendering ---")
    print(json.dumps(abstract_data, indent=2))
    print("---------------------------------------------------")


    # --- HTML Template embedded directly in the Python script ---
    html_template_content = (
"""<!DOCTYPE html>
<html lang=\"en\">
<head>
    <meta charset=\"UTF-8\">
    <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">
    <title>Graphical Abstract</title>
    <script src=\"https://cdn.tailwindcss.com\"></script>
    <!-- New font for main title -->
    <link href=\"https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Playfair+Display:wght@700&display=swap\" rel=\"stylesheet\">
    <!-- Chart.js CDN for dynamic charts -->
    <script src=\"https://cdn.jsdelivr.net/npm/chart.js\"></script>
    <style>
        body {
            font-family: 'Inter', sans-serif;
            background-color: #f0f4f8;
            display: flex;
            flex-direction: column; /* Changed to column to stack content */
            align-items: center; /* Align content to the left */
            min-height: 100vh;
            padding: 0;
            margin: 0;
            width: 100%; /* Ensure body takes full width */
        }
        .abstract-container {
            background-color: #ffffff;
            border-radius: 1.5rem;
            box-shadow: 0 10px 25px rgba(0, 0, 0, 0.1);
            padding: 2rem;
            width: calc(100% - 4rem); /* Adjust width to account for padding */
            max-width: 1600px; /* Increased max width for wider layout */
            margin: 2rem; /* Add margin around the container */
            box-sizing: border-box;
            display: flex;
            flex-direction: column;
            gap: 1.5rem;
            overflow-x: hidden;
        }
        @media (min-width: 768px) {
            .abstract-container {
                padding: 3rem;
                width: calc(100% - 6rem); /* Adjust width for larger padding */
            }
        }
        .section-title {
            font-size: 1.75rem;
            font-weight: 700;
            color: #1a202c;
            text-align: center;
            margin-bottom: 1rem;
        }
        .subsection-title {
            font-size: 1.25rem;
            font-weight: 600;
            color: #2d3748;
            margin-bottom: 0.75rem;
        }

        .hoverable-card {
            transition: transform 0.3s cubic-bezier(0.2, 0.8, 0.2, 1), box-shadow 0.3s cubic-bezier(0.2, 0.8, 0.2, 1);
        }

        .hoverable-card:hover {
            transform: translateY(-8px);
            box-shadow: 0 18px 40px rgba(0, 0, 0, 0.18);
        }
        .main-title-container h1 {
            font-family: 'Playfair Display', serif;
        }

        .grid-item {
            display: flex;
            flex-direction: column;
            align-items: center;
            text-align: center;
            padding: 1rem;
            background-color: #e2e8f0;
            border-radius: 0.75rem;
            transition: transform 0.3s cubic-bezier(0.2, 0.8, 0.2, 1), box-shadow 0.3s cubic-bezier(0.2, 0.8, 0.2, 1);
        }
        .grid-item:hover {
            transform: translateY(-8px);
            box-shadow: 0 18px 40px rgba(0, 0, 0, 0.18);
        }
        .icon {
            width: 3.5rem;
            height: 3.5rem;
            margin-bottom: 0.75rem;
            /* The color is now defined within the inline SVG itself */
        }
        .skill-tag {
            background-color: #a78bfa;
            color: white;
            padding: 0.25rem 0.75rem;
            border-radius: 9999px;
            font-size: 0.875rem;
            margin: 0.25rem;
        }
        .alert-box {
            background-color: #fbd38d;
            border-left: 4px solid #dd6b20;
            padding: 1rem;
            border-radius: 0.5rem;
            display: flex;
            align-items: center;
            gap: 0.75rem;
            font-weight: 500;
            color: #2d3748;
        }
        .metric-box {
            background-color: #63b3ed;
            color: white;
            padding: 1rem;
            border-radius: 0.75rem;
            text-align: center;
            font-weight: 600;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            gap: 0.5rem;
            transition: transform 0.3s ease-in-out, box-shadow 0.3s ease-in-out, background-color 0.3s ease-in-out;
        }
        .metric-box:hover {
            transform: translateY(-5px);
            box-shadow: 0 10px 20px rgba(0, 0, 0, 0.12);
            background-color: #4299e1;
        }
        .metric-value {
            font-size: 2.25rem;
            font-weight: 700;
            line-height: 1;
        }
        .award-box {
            background-color: #fff;
            border-radius: 0.75rem;
            padding: 1rem;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 0.5rem;
            text-align: center;
        }
        .award-rank {
            font-size: 2.5rem;
            font-weight: 800;
            color: #4c51bf;
            line-height: 1;
        }
        .award-description {
            font-size: 1rem;
            color: #4a5568;
            font-weight: 500;
        }
        /* Chart specific styling */
        .chart-container {
            position: relative;
            height: 300px;
            width: 100%;
            max-width: 600px;
            margin: 1rem auto;
        }
    </style>
</head>
<body>
    <div class=\"abstract-container\">
        <div class="mb-6 bg-blue-50 p-6 rounded-xl shadow-md text-center hoverable-card main-title-container">
            <h1 class=\"text-4xl font-extrabold text-gray-900\">
                {{ abstract.main_title }}
            </h1>
        </div>

        {% if abstract.about_section and abstract.about_section.summary.strip() %}
        <section class="mb-6 bg-white p-6 rounded-xl shadow-lg hoverable-card">
            <h2 class=\"section-title mb-4\">
                {{ abstract.about_section.title }}
            </h2>
            <div class="text-center text-gray-700">
                <p>{{ abstract.about_section.summary }}</p>
            </div>
        </section>
        {% endif %}

        <div class="grid grid-cols-1 lg:grid-cols-2 gap-8 mb-6">
            {% if abstract.metrics_section and abstract.metrics_section.metrics %}
            <section class="flex flex-col bg-white rounded-xl shadow-lg p-6 hoverable-card">
                <h2 class=\"section-title\">{{ abstract.metrics_section.title }}</h2>
                <div class=\"grid grid-cols-1 sm:grid-cols-2 gap-4 mb-4\">
                    {% for metric in abstract.metrics_section.metrics %}
                    <div class=\"metric-box\">
                        <div class=\"metric-value\">{{ metric.value }}</div>
                        <div>{{ metric.label }}</div>
                    </div>
                    {% endfor %}
                </div>
                {% if abstract.metrics_section.chart_js_data and abstract.metrics_section.chart_js_data.data and abstract.metrics_section.chart_js_data.data.datasets and abstract.metrics_section.chart_js_data.data.datasets[0].data %}
                <h3 class=\"subsection-title text-center mt-6\">{{ abstract.metrics_section.chart_title }}</h3>
                <p class=\"text-center text-gray-600 mb-4\">{{ abstract.metrics_section.chart_description }}</p>
                <div class=\"chart-container\">
                    <canvas id=\"primaryMetricsChart\"></canvas>
                </div>
                <script>
                    document.addEventListener('DOMContentLoaded', function() {
                        const ctx = document.getElementById('primaryMetricsChart');
                        if (ctx) {
                            new Chart(ctx, {{ abstract.metrics_section.chart_js_data | tojson }});
                        }
                    });
                </script>
                {% endif %}
            </section>
            {% endif %}

            {% if abstract.secondary_metrics_section and abstract.secondary_metrics_section.metrics %}
            <section class="flex flex-col bg-white rounded-xl shadow-lg p-6 hoverable-card">
                <h2 class=\"section-title\">{{ abstract.secondary_metrics_section.title }}</h2>
                <div class=\"grid grid-cols-1 sm:grid-cols-2 gap-4 mb-4\">
                    {% for metric in abstract.secondary_metrics_section.metrics %}
                    <div class=\"metric-box\">
                        <div class=\"metric-value\">{{ metric.value }}</div>
                        <div>{{ metric.label }}</div>
                    </div>
                    {% endfor %}
                </div>
                {% if abstract.secondary_metrics_section.chart_js_data and abstract.secondary_metrics_section.chart_js_data.data and abstract.secondary_metrics_section.chart_js_data.data.datasets and abstract.secondary_metrics_section.chart_js_data.data.datasets[0].data %}
                <h3 class=\"subsection-title text-center mt-6\">{{ abstract.secondary_metrics_section.chart_title }}</h3>
                <p class=\"text-center text-gray-600 mb-4\">{{ abstract.secondary_metrics_section.chart_description }}</p>
                <div class=\"chart-container\">
                    <canvas id=\"secondaryMetricsChart\"></canvas>
                </div>
                <script>
                    document.addEventListener('DOMContentLoaded', function() {
                        const ctx = document.getElementById('secondaryMetricsChart');
                        if (ctx) {
                            new Chart(ctx, {{ abstract.secondary_metrics_section.chart_js_data | tojson }});
                        }
                    });
                </script>
                {% endif %}
            </section>
            {% endif %}
        </div>
        
        {% if abstract.offerings_section and abstract.offerings_section.items %}
        <section class=\"mb-6\">
            <h2 class=\"section-title\">{{ abstract.offerings_section.title }}</h2>
            <div class=\"grid grid-cols-1 sm:grid-cols-2 md:grid-cols-3 lg:grid-cols-4 gap-4\">
                {% for item in abstract.offerings_section["items"] %}
                <div class=\"grid-item\">
                    <!-- Render dynamic SVG directly from the 'illustration_svg' field -->
                    {{ item.illustration_svg | safe }}
                    <div class=\"subsection-title\">{{ item.title }}</div>
                    <p>{{ item.description }}</p>
                </div>
                {% endfor %}
            </div>
        </section>
        {% endif %}

        {# Conclusion and Insights & Implications Sections side-by-side #}
        <div class="grid grid-cols-1 lg:grid-cols-2 gap-8 mb-6">
            {% if abstract.conclusion_section and abstract.conclusion_section.title %}
            <section class="flex flex-col bg-white rounded-xl shadow-lg p-6 hoverable-card">
                <h2 class=\"section-title\">{{ abstract.conclusion_section.title }}</h2>
                {% if abstract.conclusion_section.understanding_section and abstract.conclusion_section.understanding_section.title.strip() %}
                <div class=\"p-4 bg-gray-50 rounded shadow mb-4 flex-grow hoverable-card\">
                    <h3 class=\"subsection-title\">{{ abstract.conclusion_section.understanding_section.title }}</h3>
                    <p>{{ abstract.conclusion_section.understanding_section.description }}</p>
                    {% if abstract.conclusion_section.understanding_section.key_concepts %}
                    <div class=\"flex flex-wrap mt-2\">
                        {% for concept in abstract.conclusion_section.understanding_section.key_concepts %}
                        <span class=\"skill-tag\">{{ concept }}</span>
                        {% endfor %}
                    </div>
                    {% endif %}
                </div>
                {% endif %}
                {% if abstract.conclusion_section.outlook_section and abstract.conclusion_section.outlook_section.title.strip() %}
                <div class=\"p-4 bg-gray-50 rounded shadow flex-grow hoverable-card\">
                    <h3 class=\"subsection-title\">{{ abstract.conclusion_section.outlook_section.title }}</h3>
                    {% if abstract.conclusion_section.outlook_section.info_box_text %}
                    <div class=\"alert-box mb-2\">{{ abstract.conclusion_section.outlook_section.info_box_text }}</div>
                    {% endif %}
                    <p>{{ abstract.conclusion_section.outlook_section.description }}</p>
                </div>
                {% endif %}
            </section>
            {% endif %}

            {% if abstract.insights_and_implications_section and abstract.insights_and_implications_section.title.strip() %}
            <section class="flex flex-col bg-white rounded-xl shadow-lg p-6 hoverable-card">
                <h2 class=\"section-title\">{{ abstract.insights_and_implications_section.title }}</h2>
                {% if abstract.insights_and_implications_section.primary_insights and abstract.insights_and_implications_section.primary_insights.title.strip() %}
                <div class=\"p-4 bg-indigo-100 rounded shadow mb-4 flex-grow hoverable-card\">
                    <h3 class=\"subsection-title\">{{ abstract.insights_and_implications_section.primary_insights.title }}</h3>
                    <p>{{ abstract.insights_and_implications_section.primary_insights.description }}</p>
                </div>
                {% endif %}
                {% if abstract.insights_and_implications_section.practical_implications and abstract.insights_and_implications_section.practical_implications.title.strip() %}
                <div class=\"p-4 bg-indigo-100 rounded shadow mb-4 flex-grow hoverable-card\">
                    <h3 class=\"subsection-title\">{{ abstract.insights_and_implications_section.practical_implications.title }}</h3>
                    <p>{{ abstract.insights_and_implications_section.practical_implications.description }}</p>
                </div>
                {% endif %}
                {% if abstract.insights_and_implications_section.broader_significance and abstract.insights_and_implications_section.broader_significance.title.strip() %}
                <div class=\"p-4 bg-indigo-100 rounded shadow flex-grow hoverable-card\">
                    <h3 class=\"subsection-title\">{{ abstract.insights_and_implications_section.broader_significance.title }}</h3>
                    <p>{{ abstract.insights_and_implications_section.broader_significance.description }}</p>
                </div>
                {% endif %}
            </section>
            {% endif %}
        </div>

        <footer class=\"text-center text-sm text-gray-600 mt-4\">
            {{ abstract.footer_text }}
        </footer>
    </div>
</body>
</html>"""
)



    # Load the Jinja2 template from the string
    template = Template(html_template_content)
    # Pass JavaScript content variables to the Jinja2 rendering context (though not used in this version)
    render_context = {
        "abstract": abstract_data,
        # "js_llm_agents_section_content": js_llm_agents_section_content, # Removed as no longer passed directly
        # "js_full_abstract_content": js_full_abstract_content # Removed as no longer passed directly
    }
    rendered_html = template.render(render_context)

    # Returning HTML string for Streamlit to render
    return rendered_html


# --- Streamlit UI ---
# Set the page configuration using the layout from kal.py
st.set_page_config(
    page_title="Scientific Illustration Generator",
    layout="wide", # This ensures a wide layout
    initial_sidebar_state="expanded"
)

# Conditional imports and warnings for Streamlit UI
if Document is None:
    st.warning("python-docx library not found. .docx file support will be limited.")
if BeautifulSoup is None:
    st.warning("BeautifulSoup library not found. XML/HTML parsing might be less robust.")


# --- Logo Inclusion ---
# Set your logo file path here. It should be relative to your Streamlit app script,
# or a full URL to an image hosted online.
LOGO_PATH = "image.jpg" # Default placeholder logo
# You can change LOGO_PATH to your actual image file, e.g., "my_company_logo.png"
# For example: LOGO_PATH = "images/my_company_logo.png" if you have an 'images' folder.

# Use a columns layout for the logo to ensure it's left-aligned and not cut off
st.image("image.jpg", width=260)
st.title("Scientific Illustration Generator") # Removed icon from title to streamline

st.markdown("""
Upload your document (PDF, PPTX, TXT, XML)
and let the AI analyze its content to generate a visually appealing graphical abstract.
""")

# --- Sample Document Paths ---
# IMPORTANT: Update these paths to point to your actual sample documents
# For example: SAMPLE_DOC_DIR = "C:/Users/YourUser/Documents/sample_files"
# Or if in the same directory: SAMPLE_DOC_DIR = "."

SAMPLE_DOC_DIR = "SAMPLE_DOCUMENTS" # Assuming a 'sample_documents' folder in the same directory as this script

SAMPLE_FILES = {
    "XML": {"path": os.path.join(SAMPLE_DOC_DIR, "BR0108144.xml"), "mime": "text/xml"},
    "PPTX": {"path": os.path.join(SAMPLE_DOC_DIR, "EI__design.pptx"), "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    "PDF": {"path": os.path.join(SAMPLE_DOC_DIR, "Ramesh Tunga - AI Dev Day - Acuity and Agents.pdf"), "mime": "application/pdf"},
    #"DOCX": {"path": os.path.join(SAMPLE_DOC_DIR, "sample.docx"), "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
    "TXT": {"path": os.path.join(SAMPLE_DOC_DIR, "Strategic AI Transformation.txt"), "mime": "text/plain"},
    
}

# Use a container for controlling the main content block's width and alignment
with st.container():
    # Apply custom CSS to control the overall width of the Streamlit content (not the HTML preview iframe)
    st.markdown("""
        <style>
            /* Make Streamlit's main content area span full width and remove default padding */
            .main > div {
                padding-left: 0rem !important;
                padding-right: 0rem !important;
                padding-top: 0rem !important; /* Adjust top padding if needed */
                max-width: 100% !important;
                width: 100% !important;
            }

            /* Ensure the overall app body has no unwanted margins/padding */
            body {
                padding: 0 !important;
                margin: 0 !important;
                width: 100vw; /* Ensure body takes full viewport width */
            }

            /* Streamlit's specific div for the main content block, this also often adds padding */
            .st-emotion-cache-z5xwcq { /* Replace with actual class name if different in your Streamlit version */
                padding-left: 0rem !important;
                padding-right: 0rem !important;
            }
            
            /* Custom button styles to make them span available width better */
            .stButton > button {
                width: auto; /* Allow button to size naturally */
                min-width: 120px; /* Ensure a minimum width for download buttons */
                padding: 10px 15px; /* Slightly smaller padding for side-by-side */
                background-color: #4CAF50;
                color: white;
                font-size: 14px; /* Smaller font size for side-by-side */
                border-radius: 8px;
                border: none;
                cursor: pointer;
                transition: background-color 0.3s ease;
                box-shadow: 0 4px 6px rgba(0,0,0,0.1);
                margin-right: 0px; /* Crucial: Remove right margin */
            }

            .stButton > button:hover {
                background-color: #45a049;
            }

            /* Download button style */
            .stDownloadButton > button {
                padding: 10px 15px; /* Consistent padding with other buttons */
                background-color: #1a73e8;
                color: white;
                font-size: 14px; /* Consistent font size */
                border-radius: 8px;
                border: none;
                cursor: pointer;
                margin-top: 0px; /* Remove extra top margin for side-by-side */
                margin-right: 0px; /* Crucial: Remove right margin */
                transition: background-color 0.3s ease;
                box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            }

            .stDownloadButton > button:hover {
                background-color: #0d47a1;
            }

            /* Flex container for side-by-side download buttons */
            .download-buttons-container {
                display: flex;
                flex-wrap: wrap; /* Allow wrapping on smaller screens */
                gap: 0px; /* Changed from 10px to 0px to remove gaps */
                margin-top: 1rem;
                margin-bottom: 1rem; /* Add some space below the row of buttons */
                justify-content: flex-start; /* Align buttons to the left */
            }
            /* Styling for the "Generate Graphical Abstract" button */
            .stButton.generate-button > button {
                width: 100%; /* Make generate button full width */
                min-width: unset; /* Remove min-width constraint */
                font-size: 16px; /* Larger font for main action */
                padding: 10px 20px; /* Larger padding for main action */
            }
            /* Styling for the "View HTML in New Tab" link button */
            .button-link {
                display: inline-block;
                padding: 10px 20px;
                background-color: #007bff; /* A different color for distinction */
                color: white;
                text-align: center;
                text-decoration: none;
                font-size: 16px;
                border-radius: 8px;
                cursor: pointer;
                transition: background-color 0.3s ease;
                margin-top: 20px;
            }
            .button-link:hover {
                background-color: #0056b3;
                color: white !important;
            }
        </style>
        """, unsafe_allow_html=True)
    
    st.subheader("Download Sample Documents:")
    st.markdown('<div class="download-buttons-container">', unsafe_allow_html=True)

    # Function to create a download button for sample files
    def create_sample_download_button(file_type_label, file_path, mime_type, key):
        if os.path.exists(file_path):
            with open(file_path, "rb") as file:
                st.download_button(
                    label=f"Download sample {file_type_label}",
                    data=file.read(),
                    file_name=os.path.basename(file_path),
                    mime=mime_type,
                    key=key,
                    help=f"Click to download a sample {file_type_label} document."
                )
        else:
            st.warning(f"Sample {file_type_label} file not found at '{file_path}'. Please update SAMPLE_DOC_DIR.")

    # Create download buttons for each sample file directly within the flex container
    for file_type_label, details in SAMPLE_FILES.items():
        create_sample_download_button(file_type_label, details["path"], details["mime"], f"download_sample_{file_type_label}")
        
    st.markdown('</div>', unsafe_allow_html=True) # Close the flex container

    st.markdown("---") # Separator

    uploaded_file = st.file_uploader(
        "Upload the document",
        type=["pdf", "docx", "pptx", "txt", "xml"],
        help="Supported formats: PDF, PPTX, TXT, XML"
    )

    # Store uploaded file content in session state
    if 'uploaded_file_content' not in st.session_state:
        st.session_state['uploaded_file_content'] = None
    if 'uploaded_file_name' not in st.session_state:
        st.session_state['uploaded_file_name'] = None
    if 'uploaded_file_type' not in st.session_state:
        st.session_state['uploaded_file_type'] = None


    if uploaded_file is not None:
        st.info(f"File '{uploaded_file.name}' uploaded successfully.")
        
        # Store the uploaded file's content and metadata
        st.session_state['uploaded_file_content'] = uploaded_file.getvalue()
        st.session_state['uploaded_file_name'] = uploaded_file.name
        st.session_state['uploaded_file_type'] = uploaded_file.type

        # Reset file pointer for further processing if needed
        uploaded_file.seek(0)

        # Use a custom class for the generate button to control its width
        if st.button("Generate Graphical Abstract", key="generate_abstract_button"):
            with st.spinner("Extracting text and generating abstract... This may take a moment."):
                document_content = extract_text_from_document(uploaded_file)

                if document_content:
                    # st.subheader("Extracted Document Content Preview (first 500 chars):")
                    # st.code(document_content[:500] + "..." if len(document_content) > 500 else document_content)

                    generated_html = generate_graphical_abstract_html(document_content)

                    if generated_html:
                        st.success("Graphical Abstract generated successfully!")
                        
                        st.subheader("Generated Graphical Abstract Preview:")
                        # Display the HTML directly in Streamlit
                        st.components.v1.html(generated_html, height=800, width=None, scrolling=True) # Changed width to None

                        # Download HTML button
                        st.download_button(
                            label="Download Generated HTML",
                            data=generated_html.encode("utf-8"),
                            file_name="graphical_abstract.html",
                            mime="text/html",
                            key="download_generated_html_button", # Added a unique key for the button
                            help="Click to download the generated graphical abstract as an HTML file."
                        )

                        # Button to view HTML in a new tab using a data URI
                        b64_html = base64.urlsafe_b64encode(generated_html.encode("utf-8")).decode("utf-8")
                        data_uri = f"data:text/html;base64,{b64_html}"
                        
                        # st.markdown(
                        #     f"<a href='{data_uri}' target='_blank' class='button-link'>View HTML in New Tab</a>",
                        #     unsafe_allow_html=True
                        # )

                    else:
                        st.error("Failed to generate graphical abstract HTML. Please check the LLM response in the console for more details.")
                else:
                    st.warning("Could not extract content from the uploaded document. Please check the file and ensure it is not empty or corrupted or malformed.")
    else:
        st.info("Please upload a document to begin.")

# Removed the section for downloading original document types after upload
# as the user requested sample downloads before upload.
